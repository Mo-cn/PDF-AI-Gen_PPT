"""
题目答案PPT生成模块 - 纯模板方式 + 动画效果
"""
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from lxml import etree
from tqdm import tqdm

from .config import settings
from .models import QuestionSet


class TemplatePPTGenerator:
    def __init__(self):
        self.primary_color = RGBColor(0x3B, 0x82, 0xF6)
        self.secondary_color = RGBColor(0x63, 0x6F, 0x81)
        self.accent_color = RGBColor(0x8B, 0x5C, 0xF6)
        self.success_color = RGBColor(0x10, 0xB9, 0x81)
        self.warning_color = RGBColor(0xF5, 0x92, 0x00)
        self.text_primary = RGBColor(0x1F, 0x29, 0x37)
        self.text_secondary = RGBColor(0x64, 0x74, 0x8B)
        self.bg_light = RGBColor(0xF8, 0xFA, 0xFC)
        self.bg_card = RGBColor(0xFF, 0xFF, 0xFF)
        self.border_color = RGBColor(0xE2, 0xE8, 0xF0)
    
    def add_click_animation(self, slide, shape):
        try:
            spTree = slide.shapes._spTree
            shape_id = shape.shape_id
            
            timing_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            
            timing_xml = f'''
            <p:timing xmlns:p="{timing_ns}">
                <p:tnLst>
                    <p:par>
                        <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                            <p:childTnLst>
                                <p:seq concurrent="1" nextAc="seek">
                                    <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                                        <p:childTnLst>
                                            <p:par>
                                                <p:cTn id="3" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond evt="onClick" delay="0">
                                                            <p:tgtEl>
                                                                <p:spTgt spid="{shape_id}"/>
                                                            </p:tgtEl>
                                                        </p:cond>
                                                    </p:stCondLst>
                                                    <p:childTnLst>
                                                        <p:par>
                                                            <p:cTn id="4" fill="hold">
                                                                <p:stCondLst>
                                                                    <p:cond delay="0"/>
                                                                </p:stCondLst>
                                                                <p:childTnLst>
                                                                    <p:par>
                                                                        <p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                                                                            <p:stCondLst>
                                                                                <p:cond delay="0"/>
                                                                            </p:stCondLst>
                                                                            <p:childTnLst>
                                                                                <p:set>
                                                                                    <p:cBhvr>
                                                                                        <p:cTn id="6" dur="1" fill="hold">
                                                                                            <p:stCondLst>
                                                                                                <p:cond delay="0"/>
                                                                                            </p:stCondLst>
                                                                                        </p:cTn>
                                                                                        <p:tgtEl>
                                                                                            <p:spTgt spid="{shape_id}"/>
                                                                                        </p:tgtEl>
                                                                                        <p:attrNameLst>
                                                                                            <p:attrName>style.visibility</p:attrName>
                                                                                        </p:attrNameLst>
                                                                                    </p:cBhvr>
                                                                                    <p:to>
                                                                                        <p:strVal val="visible"/>
                                                                                    </p:to>
                                                                                </p:set>
                                                                            </p:childTnLst>
                                                                        </p:cTn>
                                                                    </p:par>
                                                                </p:childTnLst>
                                                            </p:cTn>
                                                        </p:par>
                                                    </p:childTnLst>
                                                </p:cTn>
                                            </p:par>
                                        </p:childTnLst>
                                    </p:cTn>
                                </p:seq>
                            </p:childTnLst>
                        </p:cTn>
                    </p:par>
                </p:tnLst>
            </p:timing>
            '''
        except:
            pass
    
    def create_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.bg_light
        bg.line.fill.background()
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(2),
            Inches(8), Inches(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.bg_card
        card.line.color.rgb = self.border_color
        card.line.width = Pt(1)
        
        left_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(2),
            Inches(0.15), Inches(3.5)
        )
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = self.primary_color
        left_accent.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.text_primary
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = self.text_secondary
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(12)
    
    def create_question_slide(self, prs: Presentation, question_num: int, question_data: dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.bg_light
        bg.line.fill.background()
        
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.primary_color
        header.line.fill.background()
        
        num_badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.3), Inches(0.25),
            Inches(0.7), Inches(0.7)
        )
        num_badge.fill.solid()
        num_badge.fill.fore_color.rgb = self.bg_card
        num_badge.line.fill.background()
        
        num_text = slide.shapes.add_textbox(Inches(0.3), Inches(0.35), Inches(0.7), Inches(0.5))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(question_num)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.CENTER
        
        header_text = slide.shapes.add_textbox(Inches(1.2), Inches(0.35), Inches(8), Inches(0.5))
        tf = header_text.text_frame
        p = tf.paragraphs[0]
        p.text = "选择题"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        content_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(1.5),
            Inches(9.4), Inches(5.7)
        )
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = self.bg_card
        content_card.line.color.rgb = self.border_color
        content_card.line.width = Pt(1)
        
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(1.3))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = question_data.get("content", "")
        p.font.size = Pt(18)
        p.font.color.rgb = self.text_primary
        p.font.bold = True
        
        options = question_data.get("options", [])
        option_labels = ["A", "B", "C", "D"]
        
        for i, (label, option) in enumerate(zip(option_labels, options)):
            top = Inches(3.1 + i * 0.85)
            
            option_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.6), top,
                Inches(8.8), Inches(0.7)
            )
            option_bg.fill.solid()
            option_bg.fill.fore_color.rgb = self.bg_light
            option_bg.line.color.rgb = self.border_color
            option_bg.line.width = Pt(1)
            
            label_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.75), top + Inches(0.12),
                Inches(0.45), Inches(0.45)
            )
            label_circle.fill.solid()
            label_circle.fill.fore_color.rgb = self.secondary_color
            label_circle.line.fill.background()
            
            label_text = slide.shapes.add_textbox(Inches(0.75), top + Inches(0.18), Inches(0.45), Inches(0.35))
            tf = label_text.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
            
            option_text = slide.shapes.add_textbox(Inches(1.4), top + Inches(0.15), Inches(7.8), Inches(0.5))
            tf = option_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = option
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_primary
        
        knowledge = question_data.get("knowledge_points", [])
        if knowledge:
            kb_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.3))
            tf = kb_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"📚 知识点: {', '.join(knowledge)}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.text_secondary
    
    def create_answer_slide(self, prs: Presentation, question_num: int, question_data: dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.bg_light
        bg.line.fill.background()
        
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.success_color
        header.line.fill.background()
        
        num_badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.3), Inches(0.25),
            Inches(0.7), Inches(0.7)
        )
        num_badge.fill.solid()
        num_badge.fill.fore_color.rgb = self.bg_card
        num_badge.line.fill.background()
        
        num_text = slide.shapes.add_textbox(Inches(0.3), Inches(0.35), Inches(0.7), Inches(0.5))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(question_num)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.success_color
        p.alignment = PP_ALIGN.CENTER
        
        header_text = slide.shapes.add_textbox(Inches(1.2), Inches(0.35), Inches(8), Inches(0.5))
        tf = header_text.text_frame
        p = tf.paragraphs[0]
        p.text = "答案解析"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        content_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(1.5),
            Inches(9.4), Inches(5.7)
        )
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = self.bg_card
        content_card.line.color.rgb = self.border_color
        content_card.line.width = Pt(1)
        
        correct_answer = question_data.get("correct_answer", "")
        options = question_data.get("options", [])
        option_labels = ["A", "B", "C", "D"]
        
        answer_idx = option_labels.index(correct_answer) if correct_answer in option_labels else 0
        correct_option = options[answer_idx] if answer_idx < len(options) else ""
        
        answer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), Inches(1.8),
            Inches(8.8), Inches(1.2)
        )
        answer_bg.fill.solid()
        answer_bg.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
        answer_bg.line.color.rgb = self.success_color
        answer_bg.line.width = Pt(2)
        
        answer_label = slide.shapes.add_textbox(Inches(0.9), Inches(1.95), Inches(1), Inches(0.4))
        tf = answer_label.text_frame
        p = tf.paragraphs[0]
        p.text = "正确答案"
        p.font.size = Pt(14)
        p.font.color.rgb = self.success_color
        p.font.bold = True
        
        answer_text = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(8), Inches(0.5))
        tf = answer_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"{correct_answer}. {correct_option}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.text_primary
        
        explanation = question_data.get("explanation", "")
        if explanation:
            exp_label = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(8.8), Inches(0.4))
            tf = exp_label.text_frame
            p = tf.paragraphs[0]
            p.text = "📝 解析"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.text_primary
            
            exp_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(8.8), Inches(2.5))
            tf = exp_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = explanation
            p.font.size = Pt(15)
            p.font.color.rgb = self.text_secondary
        
        knowledge = question_data.get("knowledge_points", [])
        if knowledge:
            kb_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.3))
            tf = kb_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"📚 涉及知识点: {', '.join(knowledge)}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.accent_color
    
    def generate_for_section(self, question_set: QuestionSet, output_dir: str) -> str:
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        def sanitize_filename(text):
            return "".join(c for c in text if c.isalnum() or c in (' ', '-', '_') or '\u4e00' <= c <= '\u9fff').strip()
        
        if question_set.parent_title:
            safe_parent = sanitize_filename(question_set.parent_title)
            safe_title = sanitize_filename(question_set.section_title)
            filename = f"{safe_parent}_{safe_title}"[:80]
        else:
            safe_title = sanitize_filename(question_set.section_title)
            filename = safe_title[:50] if safe_title else question_set.section_id[:8]
        
        file_path = output_path / f"{filename}.pptx"
        
        if file_path.exists():
            try:
                file_path.unlink()
            except PermissionError:
                raise PermissionError(f"文件被占用，请关闭 PowerPoint 后重试: {file_path}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        if question_set.parent_title:
            slide_subtitle = f"{question_set.parent_title} | 共 {question_set.total_count} 题"
        else:
            slide_subtitle = f"共 {question_set.total_count} 题"
        self.create_title_slide(prs, question_set.section_title, slide_subtitle)
        
        for i, q in enumerate(question_set.questions, 1):
            question_data = {
                "content": q.content,
                "options": q.options,
                "correct_answer": q.correct_answer,
                "explanation": q.explanation,
                "knowledge_points": q.knowledge_points
            }
            
            self.create_question_slide(prs, i, question_data)
            self.create_answer_slide(prs, i, question_data)
        
        prs.save(str(file_path))
        return str(file_path)
    
    def generate_all(
        self,
        question_sets: List[QuestionSet],
        output_dir: str,
        start: int = 0,
        count: Optional[int] = None
    ) -> List[str]:
        if count:
            question_sets = question_sets[start:start + count]
        else:
            question_sets = question_sets[start:]
        
        results = []
        for i, qs in enumerate(tqdm(question_sets, desc="生成PPT")):
            try:
                path = self.generate_for_section(qs, output_dir)
                results.append(path)
                print(f"  [{start + i + 1}] ✓ {qs.section_title[:30]} ({qs.total_count}题)")
            except Exception as e:
                print(f"  [{start + i + 1}] ✗ {qs.section_title[:30]}: {str(e)[:50]}")
        
        return results
