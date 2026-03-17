"""
PPT自动生成模块
"""
import signal
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from tqdm import tqdm

from .models import Section, PPTDocument, PPTSlide
from .ai_client import AIService
from .config import settings

_interrupted = False


class PPTGenerator:
    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        self.ai_service = AIService()
        
        self.title_font_size = Pt(settings.PPT_TITLE_FONT_SIZE)
        self.content_font_size = Pt(settings.PPT_CONTENT_FONT_SIZE)
        self.title_color = RGBColor(0x1F, 0x49, 0x7D)
        self.content_color = RGBColor(0x33, 0x33, 0x33)
        self.accent_color = RGBColor(0x00, 0x7A, 0xCC)
    
    def _create_presentation(self) -> Presentation:
        if self.template_path and Path(self.template_path).exists():
            return Presentation(self.template_path)
        return Presentation()
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None
        if slide_layout:
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = self.title_font_size
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.title_color
                paragraph.alignment = PP_ALIGN.CENTER
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = self.content_color
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, prs: Presentation, title: str, points: List[str]):
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.title_color
        
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, point in enumerate(points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {point}"
                p.font.size = self.content_font_size
                p.font.color.rgb = self.content_color
                p.space_after = Pt(12)
                p.level = 0
    
    def _add_summary_slide(self, prs: Presentation, title: str, points: List[str]):
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.accent_color
        
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, point in enumerate(points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"✓ {point}"
                p.font.size = self.content_font_size
                p.font.color.rgb = self.content_color
                p.space_after = Pt(12)
    
    def generate_for_section(
        self,
        section: Section,
        output_dir: str,
        use_ai: bool = True
    ) -> PPTDocument:
        prs = self._create_presentation()
        
        if use_ai:
            ppt_content = self.ai_service.generate_ppt_content(
                section_title=section.title,
                content=section.content
            )
            
            title_slide = ppt_content.get("title_slide", {})
            self._add_title_slide(
                prs,
                title_slide.get("title", section.title),
                title_slide.get("subtitle", "")
            )
            
            content_slides = ppt_content.get("content_slides", [])
            for slide_data in content_slides:
                self._add_content_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("points", [])
                )
            
            summary_slide = ppt_content.get("summary_slide", {})
            self._add_summary_slide(
                prs,
                summary_slide.get("title", "本章总结"),
                summary_slide.get("points", [])
            )
        else:
            self._add_title_slide(prs, section.title)
            
            paragraphs = section.content.split('\n')
            paragraphs = [p.strip() for p in paragraphs if p.strip()]
            
            chunk_size = 5
            for i in range(0, len(paragraphs), chunk_size):
                chunk = paragraphs[i:i+chunk_size]
                page_num = i // chunk_size + 1
                self._add_content_slide(
                    prs,
                    f"{section.title} - 第{page_num}页",
                    chunk
                )
            
            self._add_summary_slide(prs, "本章总结", paragraphs[:5])
        
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        safe_title = "".join(c for c in section.title if c.isalnum() or c in (' ', '-', '_')).strip()
        if not safe_title:
            safe_title = section.id
        file_path = output_path / f"{safe_title}.pptx"
        
        prs.save(str(file_path))
        
        return PPTDocument(
            section_id=section.id,
            section_title=section.title,
            slides=[],
            output_path=str(file_path)
        )
    
    def generate_for_all_sections(
        self,
        sections: List[Section],
        output_dir: str,
        use_ai: bool = True,
        min_content_length: int = 100
    ) -> List[PPTDocument]:
        global _interrupted
        ppt_documents = []
        
        flat_sections = self._flatten_sections(sections)
        
        valid_sections = [
            s for s in flat_sections 
            if len(s.content.strip()) >= min_content_length
        ]
        
        print(f"\n共识别 {len(flat_sections)} 个章节，其中 {len(valid_sections)} 个章节有足够内容")
        print("提示: 按 Ctrl+C 可中断\n")
        
        for i, section in enumerate(valid_sections, 1):
            if _interrupted:
                print(f"\n已中断，共生成 {len(ppt_documents)} 个PPT")
                break
            
            print(f"[{i}/{len(valid_sections)}] 正在生成PPT: {section.title[:30]}...")
            try:
                ppt_doc = self.generate_for_section(
                    section=section,
                    output_dir=output_dir,
                    use_ai=use_ai
                )
                print(f"  ✓ 完成")
                ppt_documents.append(ppt_doc)
            except KeyboardInterrupt:
                print("\n\n收到中断信号...")
                _interrupted = True
            except Exception as e:
                print(f"  ✗ 生成失败: {e}")
                continue
        
        return ppt_documents
    
    def _flatten_sections(self, sections: List[Section]) -> List[Section]:
        result = []
        for section in sections:
            result.append(section)
            if section.children:
                result.extend(self._flatten_sections(section.children))
        return result
    
    def generate_combined_ppt(
        self,
        sections: List[Section],
        output_path: str,
        use_ai: bool = True,
        min_content_length: int = 100
    ) -> str:
        prs = self._create_presentation()
        
        flat_sections = self._flatten_sections(sections)
        
        valid_sections = [
            s for s in flat_sections 
            if len(s.content.strip()) >= min_content_length
        ]
        
        print(f"\n共识别 {len(flat_sections)} 个章节，其中 {len(valid_sections)} 个章节有足够内容")
        
        for section in tqdm(valid_sections, desc="生成综合PPT"):
            if use_ai:
                ppt_content = self.ai_service.generate_ppt_content(
                    section_title=section.title,
                    content=section.content
                )
                
                title_slide = ppt_content.get("title_slide", {})
                self._add_title_slide(
                    prs,
                    title_slide.get("title", section.title),
                    title_slide.get("subtitle", "")
                )
                
                content_slides = ppt_content.get("content_slides", [])
                for slide_data in content_slides:
                    self._add_content_slide(
                        prs,
                        slide_data.get("title", ""),
                        slide_data.get("points", [])
                    )
                
                summary_slide = ppt_content.get("summary_slide", {})
                self._add_summary_slide(
                    prs,
                    summary_slide.get("title", f"{section.title} - 总结"),
                    summary_slide.get("points", [])
                )
            else:
                self._add_title_slide(prs, section.title)
                
                paragraphs = section.content.split('\n')
                paragraphs = [p.strip() for p in paragraphs if p.strip()]
                
                chunk_size = 5
                for i in range(0, len(paragraphs), chunk_size):
                    chunk = paragraphs[i:i+chunk_size]
                    page_num = i // chunk_size + 1
                    self._add_content_slide(
                        prs,
                        f"{section.title} - 第{page_num}页",
                        chunk
                    )
        
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_file))
        
        return str(output_file)
